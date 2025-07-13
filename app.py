import os
from flask import Flask, render_template, request, jsonify, session
from werkzeug.utils import secure_filename
from pptx import Presentation
import google.generativeai as genai
import textwrap

# Configure Gemini API key
genai.configure(api_key="USE_ANY_API_KEY")

# Use Gemini 2.0 Flash model
model = genai.GenerativeModel('gemini-2.0-flash')

app = Flask(__name__)
app.secret_key = os.urandom(24)
app.config['UPLOAD_FOLDER'] = 'uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Extract PPT content
def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    full_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"
    return full_text.strip()

# Clean and format paragraph output
def clean_paragraph(text):
    cleaned = text.replace("*", "").replace("\n", " ").strip()
    # Format into justified-looking text using textwrap
    return "\n".join(textwrap.fill(line, width=70) for line in cleaned.split('\n') if line.strip())

@app.route("/", methods=["GET", "POST"])
def index():
    summary = ""
    conversation = session.get("conversation", [])
    ppt_text = session.get("ppt_text", "")

    if request.method == "POST":
        ppt_file = request.files.get("ppt")
        if ppt_file:
            filename = secure_filename(ppt_file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            ppt_file.save(file_path)

            ppt_text = extract_text_from_ppt(file_path)
            session["ppt_text"] = ppt_text

            if ppt_text:
                prompt = f"Summarize the following PowerPoint content in a short paragraph:\n\n{ppt_text}"
            else:
                fallback_topic = os.path.splitext(filename)[0].replace("_", " ")
                prompt = f"No content was found in the PPT. Provide a brief informative paragraph about '{fallback_topic}' using general knowledge."

            response = model.generate_content(prompt)
            summary = clean_paragraph(response.text)
            session["summary"] = summary

    return render_template("index.html", summary=session.get("summary", ""), conversation=conversation)

@app.route("/ask", methods=["POST"])
def ask_question():
    question = request.form.get("question")
    ppt_text = session.get("ppt_text", "")

    if not question:
        return jsonify({"error": "No question provided."})

    if ppt_text:
        # Step 1: Try answering from PPT context
        prompt_ppt = f"Based only on the following PowerPoint content, answer the question below. If the answer isn't found, say 'NOT_FOUND'.\n\nQuestion: {question}\n\nContent:\n{ppt_text}"
        response_ppt = model.generate_content(prompt_ppt)
        answer = clean_paragraph(response_ppt.text.strip())

        if "NOT_FOUND" in answer or "not found" in answer.lower() or len(answer.split()) < 4:
            # Step 2: General AI search
            prompt_web = f"Answer the following question using your general knowledge from the web or outside sources:\n\nQuestion: {question}"
            response_web = model.generate_content(prompt_web)
            answer = clean_paragraph(response_web.text.strip())
    else:
        # No PPT text at all, go directly to AI
        prompt_web = f"Answer the following question using your general knowledge:\n\nQuestion: {question}"
        response_web = model.generate_content(prompt_web)
        answer = clean_paragraph(response_web.text.strip())

    # Save to chat history
    conversation = session.get("conversation", [])
    conversation.append({"question": question, "answer": answer})
    session["conversation"] = conversation

    return jsonify({"answer": answer})

# Run the app
if __name__ == "__main__":
    app.run(debug=True)
